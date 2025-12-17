from io import BytesIO
from pptx import Presentation

def generate_ppt(company, heading, date, insights):
    prs = Presentation()

    # TITLE SLIDE
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = company
    slide.placeholders[1].text = f"{heading}\nGenerated on: {date}"

    # 1 GRAPH + 1 INSIGHT PER SLIDE
    for i, insight in enumerate(insights):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = f"Graph {i+1}"
        slide.placeholders[1].text = f"Graph Placeholder\n\n{insight}"

    buffer = BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer.read()
